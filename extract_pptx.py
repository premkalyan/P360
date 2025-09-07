#!/usr/bin/env python3
"""
P360 PowerPoint Content Extractor
Extracts all text content from P360.pptx for analysis and POV creation
"""

import sys
import os
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("python-pptx not installed. Installing...")
    os.system("pip install python-pptx")
    from pptx import Presentation

def extract_pptx_content(pptx_path, output_path=None):
    """
    Extract all text content from PowerPoint presentation
    """
    if not os.path.exists(pptx_path):
        print(f"Error: File {pptx_path} not found")
        return None
    
    print(f"Extracting content from: {pptx_path}")
    
    try:
        prs = Presentation(pptx_path)
        content = []
        
        content.append("# P360 Presentation Content Extraction")
        content.append("=" * 50)
        content.append("")
        
        for i, slide in enumerate(prs.slides, 1):
            content.append(f"## Slide {i}")
            content.append("-" * 20)
            
            # Extract text from all shapes in the slide
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                content.extend(slide_text)
            else:
                content.append("(No text content)")
            
            content.append("")  # Empty line between slides
        
        # Join all content
        full_content = "\n".join(content)
        
        # Save to file if output path provided
        if output_path:
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(full_content)
            print(f"Content saved to: {output_path}")
        
        return full_content
        
    except Exception as e:
        print(f"Error extracting content: {str(e)}")
        return None

def main():
    """Main function"""
    script_dir = Path(__file__).parent
    pptx_path = script_dir / "P360.pptx"
    output_path = script_dir / "documentation" / "P360_extracted_content.md"
    
    print("P360 PowerPoint Content Extractor")
    print("=" * 40)
    
    content = extract_pptx_content(str(pptx_path), str(output_path))
    
    if content:
        print(f"\nExtraction completed successfully!")
        print(f"Content preview (first 500 chars):")
        print("-" * 40)
        print(content[:500] + "..." if len(content) > 500 else content)
        print(f"\nFull content saved to: {output_path}")
        return True
    else:
        print("Extraction failed!")
        return False

if __name__ == "__main__":
    success = main()
    sys.exit(0 if success else 1)
